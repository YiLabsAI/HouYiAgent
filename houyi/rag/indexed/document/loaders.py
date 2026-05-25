"""Indexed ingest document loaders for supported source formats."""

from __future__ import annotations

import csv
import logging
import os
import shutil
import subprocess
import uuid
from pathlib import Path

from houyi.rag.types import Document

logger = logging.getLogger(__name__)

# Canonical suffix list for ingest/document discovery in SDK and server.
SUPPORTED_DOCUMENT_SUFFIXES: tuple[str, ...] = (
    ".md",
    ".txt",
    ".rst",
    ".pdf",
    ".json",
    ".yaml",
    ".yml",
    ".csv",
    ".html",
    ".htm",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".doc",
    ".docx",
    ".pptx",
    ".epub",
)


async def load_documents(paths: list[str]) -> list[Document]:
    """Load documents from file paths.

    Supports: .txt, .md, .pdf (requires pypdf), .html/.htm, .csv,
    .xlsx/.xlsm/.xls, .doc/.docx, .pptx, .epub

    Args:
        paths: List of file paths or directories

    Returns:
        List of Document objects
    """
    documents: list[Document] = []

    for path in paths:
        path_obj = Path(path)

        if path_obj.is_dir():
            # Recursively load all files in directory
            for root, _, files in os.walk(path):
                for f in files:
                    file_path = Path(root) / f
                    doc = await _load_single_file(file_path)
                    if doc:
                        documents.append(doc)
        elif path_obj.is_file():
            doc = await _load_single_file(path_obj)
            if doc:
                documents.append(doc)

    return documents


async def _load_single_file(path: Path) -> Document | None:
    """Load a single file."""
    if not path.exists():
        return None

    suffix = path.suffix.lower()
    loader = _resolve_loader(suffix)
    if loader is None:
        return None

    try:
        return loader(path)
    except Exception:
        return None


def _resolve_loader(path_suffix: str):
    if path_suffix in {".txt", ".md", ".rst", ".json", ".yaml", ".yml"}:
        return _load_text_file
    if path_suffix == ".pdf":
        return _load_pdf_file
    if path_suffix in {".html", ".htm"}:
        return _load_html_file
    if path_suffix == ".csv":
        return _load_csv_file
    if path_suffix in {".xlsx", ".xlsm", ".xls"}:
        return _load_excel_file
    if path_suffix == ".doc":
        return _load_doc_file
    if path_suffix == ".docx":
        return _load_docx_file
    if path_suffix == ".pptx":
        return _load_pptx_file
    if path_suffix == ".epub":
        return _load_epub_file
    return None


def _load_text_file(path: Path) -> Document:
    """Load plain text file."""
    content = path.read_text(encoding="utf-8", errors="ignore")
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source=str(path),
        doc_type="text",
        metadata={"filename": path.name},
    )


def _format_tabular_rows(
    *,
    source_name: str,
    headers: list[str],
    rows: list[dict[str, str]],
    row_limit: int = 2000,
) -> str:
    """Format tabular rows into embedding-friendly text.

    We intentionally avoid markdown-table rendering for large files to reduce
    token overhead and keep each row semantically explicit.
    """
    normalized_headers = [h.strip() for h in headers if h and h.strip()]
    lines = [f"Table source: {source_name}"]
    if normalized_headers:
        lines.append(f"Columns: {', '.join(normalized_headers)}")

    if not rows:
        lines.append("No data rows.")
        return "\n".join(lines)

    total_rows = len(rows)
    keep = min(total_rows, row_limit)
    lines.append(f"Rows: {total_rows} (showing first {keep})")

    for idx, row in enumerate(rows[:row_limit], start=1):
        pairs = []
        for key in normalized_headers:
            value = str(row.get(key, "")).strip()
            if value:
                pairs.append(f"{key}: {value}")
        if not pairs:
            continue
        lines.append(f"Row {idx} - " + " | ".join(pairs))

    return "\n".join(lines)


def _load_csv_file(path: Path) -> Document | None:
    """Load CSV file.

    Uses stdlib csv by default and optionally pandas when available.
    """
    headers: list[str] = []
    rows: list[dict[str, str]] = []

    try:
        with open(path, encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.DictReader(f)
            headers = [h for h in (reader.fieldnames or []) if h]
            for row in reader:
                rows.append({k: "" if v is None else str(v) for k, v in row.items()})
    except Exception as e:
        logger.warning("Failed to parse CSV %s: %s", path, e)
        return None

    content = _format_tabular_rows(source_name=path.name, headers=headers, rows=rows)
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source=str(path),
        doc_type="csv",
        metadata={
            "filename": path.name,
            "rows": len(rows),
            "columns": len(headers),
        },
    )


def _load_doc_file(path: Path) -> Document | None:
    """Load legacy .doc files.

    This loader prefers antiword when available to avoid heavyweight runtime
    dependencies. If antiword is missing, we return None with guidance.
    """
    antiword = shutil.which("antiword")
    if antiword is None:
        logger.warning("DOC file %s skipped: install antiword to enable .doc parsing", path)
        return None

    try:
        result = subprocess.run(
            [antiword, str(path)],
            check=False,
            capture_output=True,
            text=True,
        )
    except Exception as e:
        logger.warning("Failed to parse DOC %s with antiword: %s", path, e)
        return None

    if result.returncode != 0:
        logger.warning(
            "antiword failed for %s (code=%s): %s",
            path,
            result.returncode,
            (result.stderr or "").strip(),
        )
        return None

    content = result.stdout.strip()
    if not content:
        return None

    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source=str(path),
        doc_type="doc",
        metadata={"filename": path.name},
    )


def _load_docx_file(path: Path) -> Document | None:
    """Load DOCX files via python-docx."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("DOCX file %s skipped: python-docx is not installed", path)
        return None

    try:
        doc = DocxDocument(str(path))
    except Exception as e:
        logger.warning("Failed to parse DOCX %s: %s", path, e)
        return None

    lines: list[str] = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            lines.append(text)

    # Include table content so structured docs remain searchable.
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))

    if not lines:
        return None

    return Document(
        doc_id=str(uuid.uuid4()),
        content="\n".join(lines),
        source=str(path),
        doc_type="docx",
        metadata={"filename": path.name},
    )


def _load_pptx_file(path: Path) -> Document | None:
    """Load PPTX files via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("PPTX file %s skipped: python-pptx is not installed", path)
        return None

    try:
        presentation = Presentation(str(path))
    except Exception as e:
        logger.warning("Failed to parse PPTX %s: %s", path, e)
        return None

    lines: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                stripped = text.strip()
                if stripped:
                    slide_lines.append(stripped)
        if slide_lines:
            lines.append(f"Slide {idx}")
            lines.extend(slide_lines)

    if not lines:
        return None

    return Document(
        doc_id=str(uuid.uuid4()),
        content="\n".join(lines),
        source=str(path),
        doc_type="pptx",
        metadata={"filename": path.name, "slides": len(presentation.slides)},
    )


def _load_epub_file(path: Path) -> Document | None:
    """Load EPUB files via EbookLib and BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        from ebooklib import ITEM_DOCUMENT, epub
    except ImportError:
        logger.warning("EPUB file %s skipped: install EbookLib and beautifulsoup4", path)
        return None

    try:
        book = epub.read_epub(str(path))
    except Exception as e:
        logger.warning("Failed to parse EPUB %s: %s", path, e)
        return None

    chapter_texts: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        text = soup.get_text(" ", strip=True)
        if text:
            chapter_texts.append(text)

    if not chapter_texts:
        return None

    title = ""
    title_meta = book.get_metadata("DC", "title")
    if title_meta:
        title = str(title_meta[0][0])

    return Document(
        doc_id=str(uuid.uuid4()),
        content="\n\n".join(chapter_texts),
        source=str(path),
        doc_type="epub",
        metadata={
            "filename": path.name,
            "title": title,
            "chapters": len(chapter_texts),
        },
    )


def _load_excel_file(path: Path) -> Document | None:
    """Load Excel files.

    - .xlsx/.xlsm: pandas + openpyxl
    - .xls: pandas + xlrd
    """
    try:
        import pandas as pd
    except ImportError:
        logger.warning("Excel file %s skipped: pandas is not installed", path)
        return None

    suffix = path.suffix.lower()
    engine: str | None = None
    if suffix in {".xlsx", ".xlsm"}:
        engine = "openpyxl"
    elif suffix == ".xls":
        try:
            import xlrd  # noqa: F401
        except ImportError:
            logger.warning("Excel file %s skipped: xlrd is required for .xls", path)
            return None
        engine = "xlrd"

    try:
        sheet_map = pd.read_excel(path, sheet_name=None, dtype=str, engine=engine)
    except Exception as e:
        logger.warning("Failed to parse Excel %s: %s", path, e)
        return None

    sheet_texts: list[str] = []
    total_rows = 0
    for sheet_name, frame in sheet_map.items():
        normalized = frame.fillna("").astype(str)
        headers = [str(c) for c in normalized.columns]
        rows = normalized.to_dict(orient="records")
        total_rows += len(rows)
        sheet_text = _format_tabular_rows(
            source_name=f"{path.name}::{sheet_name}",
            headers=headers,
            rows=rows,
        )
        sheet_texts.append(f"Sheet: {sheet_name}\n{sheet_text}")

    if not sheet_texts:
        return None

    return Document(
        doc_id=str(uuid.uuid4()),
        content="\n\n".join(sheet_texts),
        source=str(path),
        doc_type="excel",
        metadata={
            "filename": path.name,
            "sheets": len(sheet_map),
            "rows": total_rows,
        },
    )


def _load_pdf_file(path: Path) -> Document | None:
    """Load PDF file using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        # pypdf not installed
        return None

    reader = PdfReader(str(path))
    content = ""
    for page in reader.pages:
        content += page.extract_text() + "\n"

    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source=str(path),
        doc_type="pdf",
        metadata={
            "filename": path.name,
            "pages": len(reader.pages),
        },
    )


def _load_html_file(path: Path) -> Document:
    """Load HTML file, extracting text content."""
    content = path.read_text(encoding="utf-8", errors="ignore")

    # Simple HTML tag removal
    import re

    text = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()

    return Document(
        doc_id=str(uuid.uuid4()),
        content=text,
        source=str(path),
        doc_type="html",
        metadata={"filename": path.name},
    )
